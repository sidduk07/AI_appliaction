from pptx import Presentation

def extract_text_from_ppt(ppt_path):
    try:
        prs = Presentation(ppt_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except FileNotFoundError:
        print(f"Error: The file {ppt_path} was not found.")
    except PermissionError:
        print(f"Error: Permission denied when trying to open {ppt_path}.")
    except Exception as e:
        print(f"An unexpected error occurred: {str(e)}")
    return None